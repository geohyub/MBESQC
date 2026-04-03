"""DQR (Daily QC Report) PPT Generator.

Generates an 11-slide PowerPoint presentation — Clean Corporate style.
Designed for projection and client delivery.

  Slide 1:  Cover (split panel: dark navy / white)
  Slide 2:  Project Information (grid table)
  Slide 3:  Acquisition Settings
  Slide 4:  Line Information
  Slide 5:  Track Plot
  Slide 6-11: Surface images (Depth, StdDev, TVU, THU, Density, Backscatter)
"""

from __future__ import annotations

import datetime
from pathlib import Path

from pds_toolkit.models import GsfFile, HvfFile, PdsMetadata

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree as _etree
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ── Color Palette (GeoView Marine — matches Seismic_Processor) ─

_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
_BG         = RGBColor(0xF8, 0xFA, 0xFB)   # CLR_LIGHT_BG
_NAVY       = RGBColor(0x1E, 0x3A, 0x5F)   # CLR_NAVY (primary)
_NAVY_DARK  = RGBColor(0x00, 0x4E, 0x64)   # CLR_DARK_TEAL (gradient end)
_TEAL       = RGBColor(0x00, 0x77, 0xB6)   # CLR_TEAL (section header)
_CYAN       = RGBColor(0x00, 0xB4, 0xD8)   # CLR_CYAN (accent lines)
_ACCENT_LINE = RGBColor(0x00, 0xD4, 0xF5)  # CLR_ACCENT_LINE (thin lines)
_TEXT       = RGBColor(0x1A, 0x1A, 0x2E)   # CLR_TEXT (body)
_TEXT_SEC   = RGBColor(0x6C, 0x75, 0x7D)   # CLR_SUBTLE
_TEXT_MUTED = RGBColor(0x6C, 0x75, 0x7D)   # CLR_SUBTLE
_ROW_EVEN   = RGBColor(0xED, 0xF4, 0xF7)   # CLR_ALT_ROW
_ROW_ODD    = RGBColor(0xFF, 0xFF, 0xFF)
_BORDER     = RGBColor(0xED, 0xF4, 0xF7)   # CLR_ALT_ROW (subtle)
_KPI_BG     = RGBColor(0xE8, 0xF4, 0xF8)   # CLR_KPI_BG
_CARD_BG    = RGBColor(0x00, 0x26, 0x4D)   # CLR_CARD_BG
_ORANGE     = RGBColor(0xED, 0x89, 0x36)
_RED        = RGBColor(0xE5, 0x3E, 0x3E)

_TOTAL_SLIDES = 10  # excluding cover


def generate_dqr_ppt(
    output_path: str | Path,
    pds_meta: PdsMetadata | None = None,
    gsf_main: GsfFile | None = None,
    hvf: HvfFile | None = None,
    surface_dir: str | Path | None = None,
    project_name: str = "",
    survey_area: str = "",
    total_line_km: float = 0.0,
    qc_results: dict | None = None,
    grid_resolution: float = 1.0,
) -> None:
    if not HAS_PPTX:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Write GSF depth range as fallback for colorbar when band_stats.json is missing
    if gsf_main and gsf_main.summary and surface_dir:
        _write_gsf_depth_hint(gsf_main, Path(surface_dir))

    _slide_cover(prs, project_name, pds_meta)
    _slide_project_info(prs, pds_meta, hvf)
    _slide_acquisition(prs, pds_meta, gsf_main)
    _slide_line_info(prs, gsf_main, survey_area, total_line_km)
    _slide_track_plot(prs, surface_dir, gsf_main)

    # Pre-render all surfaces from GSF if available
    if gsf_main and gsf_main.pings and surface_dir:
        _render_all_from_gsf(gsf_main, Path(surface_dir), grid_resolution)

    surface_slides = [
        ("05", "Bathymetric Average Values Surface", "Depth"),
        ("06", "Bathymetric Standard Deviation Surface", "Std_Dev"),
        ("07", "Bathymetric TVU Surface", "TVU"),
        ("08", "Bathymetric THU Surface", "THU"),
        ("09", "Bathymetric Density Surface", "Density"),
        ("10", "Backscatter Surface", "Backscatter"),
    ]
    for num, title, surf_name in surface_slides:
        _slide_surface(prs, num, title, surf_name, surface_dir, qc_results)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ── Slide Builders ─────────────────────────────────────────────


def _slide_cover(prs, project_name, pds_meta):
    slide = _blank_slide(prs)

    # Full gradient background (Navy → DarkTeal, vertical)
    bg = slide.background.fill
    bg.gradient()
    bg.gradient_angle = 90  # degrees (top→bottom)
    bg.gradient_stops[0].color.rgb = _NAVY
    bg.gradient_stops[1].color.rgb = _NAVY_DARK

    # Left cyan stripe
    _rect(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), _CYAN)

    # "DAILY QC REPORT" subtitle
    _text(slide, "DAILY QC REPORT",
          Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.5),
          size=16, color=_CYAN, bold=False, spacing=3)

    # Main title
    _text(slide, "Multibeam Echosounder",
          Inches(0.4), Inches(1.8), Inches(10), Inches(1.2),
          size=40, color=_WHITE, bold=True)

    # Accent line
    _rect(slide, Inches(0.4), Inches(3.2), Inches(3.0), Pt(3), _ACCENT_LINE)

    # Project info card
    name = project_name or (pds_meta.project_name if pds_meta else "MBES Survey")
    if len(name) > 60:
        name = name[:57] + "..."

    vessel = pds_meta.vessel_name if pds_meta and pds_meta.vessel_name else ""
    date_str = datetime.datetime.now().strftime("%Y. %m. %d")

    card_y = Inches(3.6)
    card_items = [
        ("Project", name),
        ("Date", date_str),
    ]
    if vessel:
        card_items.append(("Vessel", vessel))

    # Card background (semi-transparent dark)
    card_h = len(card_items) * 0.42 + 0.2
    card = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(0.4), card_y, Inches(6.0), Inches(card_h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _CARD_BG
    card.line.fill.background()
    # Set alpha for semi-transparency
    spPr = card._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is not None:
        srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = _etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", "75000")  # 75% opacity

    for i, (label, value) in enumerate(card_items):
        row_y = card_y + Inches(0.1 + i * 0.42)
        _text(slide, label,
              Inches(0.7), row_y, Inches(2), Inches(0.2),
              size=8, bold=True, color=_CYAN)
        _text(slide, value,
              Inches(0.7), row_y + Inches(0.18), Inches(5), Inches(0.24),
              size=14, bold=True, color=_WHITE)

    # Company name (bottom)
    _text(slide, "GEOVIEW CO., Ltd.",
          Inches(0.4), Inches(6.6), Inches(4), Inches(0.4),
          size=12, color=RGBColor(0x6C, 0x75, 0x7D))

    # Footer
    _footer_bar(slide, prs.slide_width, page=1)


def _slide_project_info(prs, pds_meta, hvf):
    slide = _blank_slide(prs)
    _set_bg(slide, _BG)
    _header_bar(slide, "Project Information", prs.slide_width, number="01")
    _footer_bar(slide, prs.slide_width, page=2)

    rows = []
    if pds_meta:
        rows.append(("Vessel", pds_meta.vessel_name or "-"))
        rows.append(("Survey Type", pds_meta.survey_type or "-"))
        rows.append(("Coordinate System",
                      f"{pds_meta.coord_system_group} / {pds_meta.coord_system_name}"))
        rows.append(("Units", pds_meta.system_units or "-"))

    if hvf:
        for s in hvf.sensors:
            if "DepthSensor" in s.name or "Multibeam" in s.name:
                rows.append(("Transducer Offset",
                              f"X = {s.x:.3f}   Y = {s.y:.3f}   Z = {s.z:.3f}"))
                rows.append(("Mount Angles",
                              f"Pitch = {s.pitch:.3f}   Roll = {s.roll:.3f}   Heading = {s.heading:.3f}"))
                break
        sensor_names = sorted(set(_clean_sensor_name(s.name) for s in hvf.sensors))
        if sensor_names:
            if len(sensor_names) <= 8:
                rows.append(("Registered Sensors", ", ".join(sensor_names)))
            else:
                rows.append(("Registered Sensors",
                              f"{len(sensor_names)} sensors: {', '.join(sensor_names[:5])}, ..."))

    if not rows:
        rows.append(("Status", "No project metadata available"))

    _kv_table(slide, rows, Inches(0.3), Inches(0.85), Inches(12.7))


def _slide_acquisition(prs, pds_meta, gsf_main):
    slide = _blank_slide(prs)
    _set_bg(slide, _BG)
    _header_bar(slide, "Acquisition Settings", prs.slide_width, number="02")
    _footer_bar(slide, prs.slide_width, page=3)

    rows = []
    if pds_meta:
        geom = pds_meta.sections.get("GEOMETRY", {})
        for k, v in geom.items():
            if k.startswith("Offset(") and "T50" in v:
                rows.append(("MBES Model", v.split(",")[0]))

    if gsf_main and gsf_main.pings:
        p0 = gsf_main.pings[0]
        rows.append(("Number of Beams", str(p0.num_beams)))
        if p0.depth is not None:
            import numpy as np
            valid_depth = p0.depth[p0.depth != 0]
            if len(valid_depth) > 0:
                rows.append(("Swath Depth (Ping 1)",
                              f"{valid_depth.min():.1f} ~ {valid_depth.max():.1f} m"))
        if hasattr(p0, 'heading') and p0.heading:
            rows.append(("Heading (Start)", f"{p0.heading:.1f}\u00b0"))
        rows.append(("Total Pings", str(len(gsf_main.pings))))

    if gsf_main and gsf_main.summary:
        s = gsf_main.summary
        d_min = min(s.min_depth, s.max_depth)
        d_max = max(s.min_depth, s.max_depth)
        if abs(d_max - d_min) < 0.05:
            rows.append(("Depth (Survey)", f"{d_min:.1f} m"))
        else:
            rows.append(("Depth Range (Survey)", f"{d_min:.1f} ~ {d_max:.1f} m"))

    if gsf_main and gsf_main.svp_profiles:
        svp = gsf_main.svp_profiles[0]
        if svp.num_points > 0:
            rows.append(("Surface Sound Velocity", f"{svp.sound_velocity[0]:.1f} m/s"))
            rows.append(("SVP Profile Points", str(svp.num_points)))
            rows.append(("SVP Max Depth", f"{svp.depth[-1]:.1f} m"))

    if not rows:
        rows.append(("Status", "No acquisition data available"))

    _kv_table(slide, rows, Inches(0.3), Inches(0.85), Inches(12.7))


def _slide_line_info(prs, gsf_main, survey_area, total_line_km):
    slide = _blank_slide(prs)
    _set_bg(slide, _BG)
    _header_bar(slide, "Line Information", prs.slide_width, number="03")
    _footer_bar(slide, prs.slide_width, page=4)

    if total_line_km <= 0 and gsf_main and gsf_main.pings:
        try:
            import numpy as np
            lats = np.array([p.latitude for p in gsf_main.pings if p.latitude != 0])
            lons = np.array([p.longitude for p in gsf_main.pings if p.longitude != 0])
            if len(lats) > 1:
                total_line_km = _compute_line_km(lats, lons)
        except Exception:
            pass

    rows = [
        ("Survey Area", survey_area or "N/A"),
        ("Acquired Lines", f"{total_line_km:.1f} km"),
    ]
    if gsf_main and gsf_main.summary:
        s = gsf_main.summary
        d_min = min(s.min_depth, s.max_depth)
        d_max = max(s.min_depth, s.max_depth)
        if abs(d_max - d_min) < 0.05:
            rows.append(("Depth", f"{d_min:.1f} m"))
        else:
            rows.append(("Depth Range", f"{d_min:.1f} ~ {d_max:.1f} m"))
        rows.append(("Latitude", f"{s.min_latitude:.6f} ~ {s.max_latitude:.6f}"))
        rows.append(("Longitude", f"{s.min_longitude:.6f} ~ {s.max_longitude:.6f}"))
        start_t = getattr(s, "start_time", None) or getattr(s, "min_time", None)
        end_t = getattr(s, "end_time", None) or getattr(s, "max_time", None)
        if start_t and end_t:
            rows.append(("Time Span", f"{start_t} ~ {end_t}"))

    _kv_table(slide, rows, Inches(0.3), Inches(0.85), Inches(12.7))


def _slide_track_plot(prs, surface_dir, gsf_main=None):
    slide = _blank_slide(prs)
    _set_bg(slide, _BG)
    _header_bar(slide, "Track Plot", prs.slide_width, number="04")
    _footer_bar(slide, prs.slide_width, page=5)

    embedded = False

    if surface_dir:
        for name in ["trackplot", "track_plot", "Track", "navigation"]:
            for ext in [".png", ".jpg", ".jpeg"]:
                p = Path(surface_dir) / f"{name}{ext}"
                if p.exists():
                    _image_frame(slide, p,
                                 Inches(0.3), Inches(0.72), Inches(12.7), Inches(5.83))
                    embedded = True
                    break
            if embedded:
                break

    if not embedded and gsf_main and gsf_main.pings:
        result = _render_track_plot(gsf_main, Path(surface_dir) if surface_dir else None)
        if result:
            track_png, auto_km = result
            if track_png and Path(track_png).exists():
                _image_frame(slide, track_png,
                             Inches(0.3), Inches(0.72), Inches(12.7), Inches(5.83))
                embedded = True

    if not embedded:
        _placeholder(slide, "Track Plot",
                     "CARIS Export or survey software generated track image required")


def _slide_surface(prs, number, title, surf_name, surface_dir, qc_results):
    slide = _blank_slide(prs)
    _set_bg(slide, _BG)
    page = int(number) + 1
    _header_bar(slide, title, prs.slide_width, number=number)
    _footer_bar(slide, prs.slide_width, page=page)

    embedded = False
    if surface_dir:
        png_path = _render_surface_png(Path(surface_dir), surf_name)
        if png_path and png_path.exists():
            _image_frame(slide, png_path,
                         Inches(0.3), Inches(0.72), Inches(12.7), Inches(5.43))
            embedded = True

    if not embedded:
        _placeholder(slide, title,
                     "Surface image from CARIS RenderRaster required")

    # Comment line
    _text(slide, "Comment :",
          Inches(0.3), Inches(6.4), Inches(1.5), Inches(0.3),
          size=9, bold=True, color=_TEXT_SEC)
    _rect(slide, Inches(1.8), Inches(6.45), Inches(11.2), Inches(0.01), _BORDER)


# ── Layout Components ──────────────────────────────────────────


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _header_bar(slide, title, slide_width, number=None):
    """Gradient header bar (Navy→DarkTeal) with cyan accent line."""
    # Gradient header background
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), slide_width, Inches(0.58))
    fill = header.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = _NAVY
    fill.gradient_stops[1].color.rgb = _NAVY_DARK
    header.line.fill.background()

    # Cyan accent line below header
    _rect(slide, Inches(0), Inches(0.58), slide_width, Pt(2.5), _ACCENT_LINE)

    # Title (white on dark)
    _text(slide, title,
          Inches(0.35), Inches(0.05), Inches(9), Inches(0.48),
          size=17, bold=True, color=_WHITE)

    # Page counter (cyan text)
    if number:
        _text(slide, f"{number} / {_TOTAL_SLIDES:02d}",
              Inches(11.2), Inches(0.10), Inches(1.8), Inches(0.38),
              size=12, color=_CYAN, align=PP_ALIGN.RIGHT)


def _footer_bar(slide, slide_width, page=1):
    """Navy footer with accent line (matches Seismic_Processor)."""
    footer_h = Inches(0.22)
    footer_y = Inches(7.5) - footer_h
    # Accent line above footer
    _rect(slide, Inches(0), footer_y - Pt(2), slide_width, Pt(2), _ACCENT_LINE)
    # Navy footer bg
    _rect(slide, Inches(0), footer_y, slide_width, footer_h, _NAVY)
    # Company text
    _text(slide, "Geoview  |  Daily QC Report",
          Inches(0.35), footer_y + Inches(0.02), Inches(6), Inches(0.18),
          size=7, color=_WHITE)
    # Page number (cyan)
    _text(slide, f"Page {page}",
          Inches(11.5), footer_y + Inches(0.02), Inches(1.5), Inches(0.18),
          size=7, color=_CYAN, align=PP_ALIGN.RIGHT)


def _kv_table(slide, rows, left, top, width):
    """Professional grid table with proper borders and alternating rows."""
    n_rows = len(rows)
    n_cols = 2
    key_w = Inches(3.5)
    val_w = width - key_w

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.55 * n_rows))
    tbl = tbl_shape.table

    # Set column widths
    tbl.columns[0].width = key_w
    tbl.columns[1].width = val_w

    for i, (key, val) in enumerate(rows):
        row = tbl.rows[i]
        row.height = Inches(0.72) if len(val) > 80 else Inches(0.55)

        # Key cell
        key_cell = row.cells[0]
        _style_cell(key_cell, key, size=14, bold=True, color=_NAVY,
                    fill=_ROW_EVEN, align=PP_ALIGN.LEFT)
        _set_cell_borders(key_cell, _BORDER)

        # Value cell
        val_cell = row.cells[1]
        fill = _ROW_ODD if i % 2 == 0 else _ROW_EVEN
        val_size = 12 if len(val) > 80 else 13
        _style_cell(val_cell, val, size=val_size, bold=False, color=_TEXT,
                    fill=fill, align=PP_ALIGN.LEFT)
        _set_cell_borders(val_cell, _BORDER)


def _style_cell(cell, text, size=13, bold=False, color=_TEXT,
                fill=_WHITE, align=PP_ALIGN.LEFT):
    """Apply consistent formatting to a table cell."""
    cell.text = ""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = _etree.SubElement(tcPr, qn("a:solidFill"))
    _etree.SubElement(solidFill, qn("a:srgbClr")).set("val", str(fill))

    # Margins
    tcPr.set("marL", str(Inches(0.25)))
    tcPr.set("marR", str(Inches(0.15)))
    tcPr.set("marT", str(Inches(0.06)))
    tcPr.set("marB", str(Inches(0.06)))

    # Text
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Pretendard"
    p.alignment = align
    cell.text_frame.word_wrap = True


def _set_cell_borders(cell, color, width=Pt(0.5)):
    """Set all 4 borders of a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    color_hex = str(color)
    for edge in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = _etree.SubElement(tcPr, qn(edge))
        ln.set("w", str(int(width)))
        ln.set("cmpd", "sng")

        sf = _etree.SubElement(ln, qn("a:solidFill"))
        _etree.SubElement(sf, qn("a:srgbClr")).set("val", color_hex)

        _etree.SubElement(ln, qn("a:prstDash")).set("val", "solid")


def _image_frame(slide, image_path, left, top, max_w, max_h):
    """Add image preserving aspect ratio, centered within the available area."""
    from PIL import Image as PILImage
    try:
        with PILImage.open(str(image_path)) as img:
            img_w, img_h = img.size
    except Exception:
        img_w, img_h = 1600, 900

    aspect = img_w / img_h
    max_w_in = max_w / 914400
    max_h_in = max_h / 914400

    if aspect > (max_w_in / max_h_in):
        w = max_w_in
        h = w / aspect
    else:
        h = max_h_in
        w = h * aspect

    x = (left / 914400) + (max_w_in - w) / 2
    y = (top / 914400) + (max_h_in - h) / 2

    slide.shapes.add_picture(str(image_path),
                              Inches(x), Inches(y), Inches(w), Inches(h))


def _placeholder(slide, title, hint):
    """Dashed border placeholder for missing images."""
    frame = slide.shapes.add_shape(
        1,  # RECTANGLE (not rounded)
        Inches(0.3), Inches(0.72), Inches(12.7), Inches(5.83),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = _WHITE
    frame.line.color.rgb = _BORDER
    frame.line.width = Pt(1.5)
    frame.line.dash_style = 2  # DASH

    _text(slide, title,
          Inches(4), Inches(3.4), Inches(5.5), Inches(0.5),
          size=18, bold=True, color=_TEXT_MUTED, align=PP_ALIGN.CENTER)
    _text(slide, hint,
          Inches(3), Inches(4.0), Inches(7.5), Inches(0.5),
          size=13, color=_TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ── Helpers ────────────────────────────────────────────────────


def _render_all_from_gsf(gsf_main, surface_dir: Path, grid_resolution: float = 1.0):
    """Render Depth/Density/Std_Dev/TVU/THU surfaces from GSF sounding data.

    Direct gridding from beam-level data ensures image and colorbar match 100%.
    Only generates PNGs that don't already exist.
    """
    try:
        import numpy as np
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from scipy import stats as sp_stats

        # Check which surfaces need rendering
        needed = {}
        # Check which surfaces need GSF-based rendering
        # Prefer CARIS outputs (proper georef) when available:
        #   Depth/Density: CARIS TIF + ASCII bands
        #   TVU/THU: CARIS TPU bands (tpu_bands.txt + surface_tpu.tif)
        #   Std_Dev: GSF (CARIS ASCII Std_Dev is unreliable)
        _has_caris_tif = any(
            (surface_dir / f"surface_{n.lower()}.tif").exists() or
            (surface_dir / f"{n}.tif").exists()
            for n in ["Depth", "depth"]
        )
        _has_caris_tpu = (surface_dir / "tpu_bands.txt").exists() and \
                         (surface_dir / "surface_tpu.tif").exists()
        for name in ["Depth", "Density", "Std_Dev", "TVU", "THU"]:
            png = surface_dir / f"{name}.png"
            if png.exists():
                continue
            if name in ("Depth", "Density") and _has_caris_tif:
                continue  # _render_surface_png handles via CARIS TIF/ASCII
            if name in ("TVU", "THU") and _has_caris_tpu:
                continue  # _render_surface_png handles via TPU bands
            needed[name] = True
        if not needed:
            return

        # Check GSF has required data
        p0 = gsf_main.pings[0]
        if p0.depth is None or p0.across_track is None:
            return

        # Collect beam positions in UTM
        try:
            from pyproj import Transformer
            mean_lon = np.mean([p.longitude for p in gsf_main.pings[:10]
                                if p.longitude != 0])
            mean_lat = np.mean([p.latitude for p in gsf_main.pings[:10]
                                if p.latitude != 0])
            zone = int((mean_lon + 180) / 6) + 1
            epsg = 32600 + zone if mean_lat >= 0 else 32700 + zone
            if 124 < mean_lon < 132 and 33 < mean_lat < 39:
                epsg = 32652
            transformer = Transformer.from_crs("EPSG:4326", f"EPSG:{epsg}",
                                                always_xy=True)
        except Exception:
            return

        # Compute actual course from consecutive ping positions
        # (GSF heading field is often constant/wrong after CARIS export)
        _ping_e, _ping_n = [], []
        _valid_idx = []
        for i, p in enumerate(gsf_main.pings):
            if p.latitude == 0 or p.depth is None:
                continue
            e, n = transformer.transform(p.longitude, p.latitude)
            _ping_e.append(e)
            _ping_n.append(n)
            _valid_idx.append(i)
        _ping_e = np.array(_ping_e)
        _ping_n = np.array(_ping_n)

        # Course from positions spaced >= 10m apart (robust to GPS jitter)
        _MIN_BASELINE = 10.0  # meters
        _dists = np.sqrt(np.diff(_ping_e) ** 2 + np.diff(_ping_n) ** 2)
        _cum_dist = np.concatenate([[0], np.cumsum(_dists)])

        _courses = np.zeros(len(_ping_e))
        for j in range(len(_ping_e)):
            # Look ahead/behind for baseline
            fwd = min(len(_ping_e) - 1,
                      np.searchsorted(_cum_dist, _cum_dist[j] + _MIN_BASELINE))
            bwd = max(0,
                      np.searchsorted(_cum_dist, _cum_dist[j] - _MIN_BASELINE))
            if fwd == j:
                fwd = min(len(_ping_e) - 1, j + 1)
            if bwd == j:
                bwd = max(0, j - 1)
            de = _ping_e[fwd] - _ping_e[bwd]
            dn = _ping_n[fwd] - _ping_n[bwd]
            _courses[j] = np.degrees(np.arctan2(de, dn)) % 360

        # Build course lookup: valid_idx[j] → course[j]
        _course_map = dict(zip(_valid_idx, _courses))

        xs, ys, depths = [], [], []
        # Separate coords for TVU/THU (some pings may lack error arrays)
        xs_ve, ys_ve, vert_errors = [], [], []
        xs_he, ys_he, horiz_errors = [], [], []
        for i, p in enumerate(gsf_main.pings):
            if p.latitude == 0 or p.depth is None:
                continue

            # Use actual array length (CARIS export may strip rejected beams)
            n = len(p.depth)
            bx, by = transformer.transform(
                np.full(n, p.longitude), np.full(n, p.latitude))
            if p.across_track is not None and len(p.across_track) == n:
                # Use computed course (reliable) instead of GSF heading (often wrong)
                course_deg = _course_map.get(i, p.heading or 0)
                h = np.radians(course_deg)
                # Across-track: perpendicular to course (starboard positive)
                bx += p.across_track * np.cos(h)
                by -= p.across_track * np.sin(h)
                # Along-track: parallel to course (forward positive)
                if p.along_track is not None and len(p.along_track) == n:
                    bx += p.along_track * np.sin(h)
                    by += p.along_track * np.cos(h)

            # Filter rejected beams via beam_flags
            mask = np.ones(n, dtype=bool)
            if p.beam_flags is not None and len(p.beam_flags) == n:
                mask = (p.beam_flags == 0)  # 0 = accepted
            # Also filter zero/negative depths
            mask &= (p.depth > 0)

            xs.append(bx[mask])
            ys.append(by[mask])
            depths.append(p.depth[mask])
            if p.vert_error is not None and len(p.vert_error) == n:
                ve = np.clip(p.vert_error[mask], 0, 100)
                vert_errors.append(ve)
                xs_ve.append(bx[mask])
                ys_ve.append(by[mask])
            if p.horiz_error is not None and len(p.horiz_error) == n:
                he = np.clip(p.horiz_error[mask], 0, 100)
                horiz_errors.append(he)
                xs_he.append(bx[mask])
                ys_he.append(by[mask])

        if not xs:
            return

        xs = np.concatenate(xs)
        ys = np.concatenate(ys)
        depths = np.concatenate(depths)

        res = grid_resolution
        xbins = np.arange(xs.min(), xs.max() + res, res)
        ybins = np.arange(ys.min(), ys.max() + res, res)

        cell_area = f"{res:.0f}m" if res >= 1 else f"{res}m"
        _BAND_LABEL = {"Depth": "Depth (m)", "Std_Dev": "Std Dev (m)",
                       "TVU": "TVU (m)", "THU": "THU (m)",
                       "Density": f"Soundings / {cell_area}\u00b2 cell"}
        _CMAP = {"Depth": "jet", "Std_Dev": "YlOrRd", "TVU": "RdYlBu_r",
                 "THU": "RdYlBu_r", "Density": "viridis"}

        def _grid_and_render(name, values, statistic="mean",
                             xs_override=None, ys_override=None,
                             values_override=None,
                             pctl_lo=2, pctl_hi=98):
            png = surface_dir / f"{name}.png"
            if png.exists():
                return
            _xs = xs_override if xs_override is not None else xs
            _ys = ys_override if ys_override is not None else ys
            _vals = values_override if values_override is not None else values
            grid, _, _, _ = sp_stats.binned_statistic_2d(
                _xs, _ys, _vals, statistic=statistic, bins=[xbins, ybins])
            grid = grid.T
            grid[grid == 0] = np.nan

            valid = ~np.isnan(grid)
            if not valid.any():
                return

            # Crop to data extent
            rv = np.any(valid, axis=1)
            cv = np.any(valid, axis=0)
            r0, r1 = np.where(rv)[0][[0, -1]]
            c0, c1 = np.where(cv)[0][[0, -1]]
            pad = max(2, int(0.02 * max(r1 - r0, c1 - c0)))
            grid = grid[max(0, r0 - pad):r1 + pad + 1,
                         max(0, c0 - pad):c1 + pad + 1]

            # Percentile clipping for colorbar range
            valid_vals = grid[~np.isnan(grid)]
            vmin = float(np.percentile(valid_vals, pctl_lo))
            vmax = float(np.percentile(valid_vals, pctl_hi))
            if abs(vmax - vmin) < 1e-10:
                vmin, vmax = float(valid_vals.min()), float(valid_vals.max())

            # Rotate 90° if very elongated N-S → fits landscape PPT slide
            nrows, ncols = grid.shape
            rotated = False
            if nrows > ncols * 2.5:
                grid = np.rot90(grid, k=1)  # rotate CCW
                nrows, ncols = grid.shape
                rotated = True

            # Adaptive figsize: target landscape for PPT (12.7" x 5.4" frame)
            aspect = nrows / max(ncols, 1)
            fig_h = 8
            fig_w = max(10, fig_h / max(aspect, 0.3))
            fig_w = min(fig_w, 20)
            fig_w += 2  # colorbar space

            cmap = plt.get_cmap(_CMAP.get(name, "jet")).copy()
            cmap.set_bad("white")

            fig, (ax, cax) = plt.subplots(1, 2, figsize=(fig_w, fig_h),
                gridspec_kw={"width_ratios": [16, 1], "wspace": 0.05},
                facecolor="white")
            im = ax.imshow(grid, cmap=cmap, aspect="equal",
                           interpolation="nearest", vmin=vmin, vmax=vmax)
            _draw_colorbar(plt, cax, vmin, vmax, name, _BAND_LABEL,
                           mappable=im)
            ax.set_xticks([])
            ax.set_yticks([])
            for spine in ax.spines.values():
                spine.set_visible(False)

            fig.savefig(str(png), dpi=200, bbox_inches="tight",
                        facecolor="white", edgecolor="none")
            plt.close(fig)

        # Render each surface
        if "Depth" in needed:
            _grid_and_render("Depth", depths)

        if "Density" in needed:
            # Density = raw sounding count per grid cell
            _grid_and_render("Density", depths, statistic="count",
                             pctl_lo=5, pctl_hi=95)

        if "Std_Dev" in needed:
            # Std_Dev: standard deviation of ALL soundings in each cell
            # Same coverage as Depth (consistent with CARIS Std_Dev band)
            _grid_and_render("Std_Dev", depths, statistic="std",
                             pctl_lo=0, pctl_hi=95)

        if "TVU" in needed and vert_errors:
            _grid_and_render("TVU", None,
                             xs_override=np.concatenate(xs_ve),
                             ys_override=np.concatenate(ys_ve),
                             values_override=np.concatenate(vert_errors))

        if "THU" in needed and horiz_errors:
            _grid_and_render("THU", None,
                             xs_override=np.concatenate(xs_he),
                             ys_override=np.concatenate(ys_he),
                             values_override=np.concatenate(horiz_errors))

    except Exception as e:
        import logging
        logging.getLogger(__name__).warning("GSF surface render failed: %s", e)


def _render_tpu_from_gsf(gsf_main, surface_dir: Path):
    """Render TVU/THU surface images from GSF sounding-level TPU data."""
    try:
        import numpy as np
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker
        from scipy import stats as sp_stats

        # Check if GSF has TPU data
        p0 = gsf_main.pings[0]
        if not hasattr(p0, 'vert_error') or p0.vert_error is None:
            return
        if len(p0.vert_error) == 0 or p0.vert_error.max() == 0:
            return

        # Collect beam positions + errors
        xs, ys, data_map = [], [], {"TVU": [], "THU": []}
        try:
            from pyproj import Transformer
            mean_lon = np.mean([p.longitude for p in gsf_main.pings[:10]
                                if p.longitude != 0])
            mean_lat = np.mean([p.latitude for p in gsf_main.pings[:10]
                                if p.latitude != 0])
            zone = int((mean_lon + 180) / 6) + 1
            epsg = 32600 + zone if mean_lat >= 0 else 32700 + zone
            if 124 < mean_lon < 132 and 33 < mean_lat < 39:
                epsg = 32652
            transformer = Transformer.from_crs("EPSG:4326", f"EPSG:{epsg}",
                                                always_xy=True)
        except Exception:
            return

        for p in gsf_main.pings:
            if p.latitude == 0 or p.vert_error is None:
                continue
            n = p.num_beams
            bx, by = transformer.transform(
                np.full(n, p.longitude), np.full(n, p.latitude))
            if p.across_track is not None:
                h = np.radians(p.heading or 0)
                bx += p.across_track * np.cos(h)
                by -= p.across_track * np.sin(h)
                if p.along_track is not None:
                    bx += p.along_track * np.sin(h)
                    by += p.along_track * np.cos(h)
            xs.append(bx)
            ys.append(by)
            data_map["TVU"].append(p.vert_error)
            if p.horiz_error is not None:
                data_map["THU"].append(p.horiz_error)

        if not xs:
            return

        xs = np.concatenate(xs)
        ys = np.concatenate(ys)

        _BAND_LABEL = {"TVU": "TVU (m)", "THU": "THU (m)"}

        for name in ["TVU", "THU"]:
            png_path = surface_dir / f"{name}.png"
            if png_path.exists():
                continue  # Already rendered
            if not data_map[name]:
                continue

            values = np.concatenate(data_map[name])
            if values.max() == 0:
                continue

            # Bin into 1m grid
            res = 1.0
            xbins = np.arange(xs.min(), xs.max() + res, res)
            ybins = np.arange(ys.min(), ys.max() + res, res)
            grid, _, _, _ = sp_stats.binned_statistic_2d(
                xs, ys, values, statistic="mean", bins=[xbins, ybins])
            grid = grid.T
            grid[grid == 0] = np.nan

            # Crop
            valid = ~np.isnan(grid)
            if not valid.any():
                continue
            rv = np.any(valid, axis=1)
            cv = np.any(valid, axis=0)
            r0, r1 = np.where(rv)[0][[0, -1]]
            c0, c1 = np.where(cv)[0][[0, -1]]
            pad = max(2, int(0.02 * max(r1 - r0, c1 - c0)))
            grid = grid[max(0, r0 - pad):r1 + pad + 1,
                         max(0, c0 - pad):c1 + pad + 1]

            cmap = plt.get_cmap("RdYlBu_r").copy()
            cmap.set_bad("white")

            fig, (ax, cax) = plt.subplots(1, 2, figsize=(14, 10),
                gridspec_kw={"width_ratios": [16, 1], "wspace": 0.04},
                facecolor="white")
            im = ax.imshow(grid, cmap=cmap, aspect="equal",
                           interpolation="nearest")
            vmin, vmax = float(np.nanmin(grid)), float(np.nanmax(grid))
            _draw_colorbar(plt, cax, vmin, vmax, name, _BAND_LABEL,
                           mappable=im)
            ax.set_xticks([])
            ax.set_yticks([])
            for spine in ax.spines.values():
                spine.set_visible(False)

            fig.savefig(str(png_path), dpi=200, bbox_inches="tight",
                        facecolor="white", edgecolor="none")
            plt.close(fig)
    except Exception:
        pass


def _write_gsf_depth_hint(gsf_main, surface_dir: Path):
    """Write GSF depth range to band_stats.json as fallback for colorbar."""
    stats_file = surface_dir / "band_stats.json"
    if stats_file.exists():
        return  # already has real stats
    try:
        import json
        s = gsf_main.summary
        d_min = min(s.min_depth, s.max_depth)
        d_max = max(s.min_depth, s.max_depth)
        stats = {"Depth": {"min": round(d_min, 3), "max": round(d_max, 3)}}
        with open(str(stats_file), "w") as f:
            json.dump(stats, f, indent=2)
    except Exception:
        pass


def _clean_sensor_name(name: str) -> str:
    """Strip timestamps and IDs from sensor names.

    'DepthSensor_T1_2025-230 00:00:00' → 'DepthSensor (T1)'
    'DepthSensor_T1_2025-044_00:00:00' → 'DepthSensor (T1)'
    """
    import re
    # Remove date-time suffix: underscore OR space before HH:MM:SS
    cleaned = re.sub(r'_\d{4}-\d{2,3}[\s_]\d{2}:\d{2}:\d{2}$', '', name)
    # Convert DepthSensor_T1 → DepthSensor (T1)
    m = re.match(r'^(\w+Sensor)_([A-Z]\d+)$', cleaned)
    if m:
        return f"{m.group(1)} ({m.group(2)})"
    return cleaned


def _draw_colorbar(plt, cax, vmin, vmax, surf_name, band_labels, mappable=None):
    """Draw a polished colorbar on the given axes."""
    import matplotlib.cm as cm
    import matplotlib.ticker as mticker

    _CMAP = {
        "Depth": "jet", "Std_Dev": "YlOrRd", "TVU": "RdYlBu_r",
        "THU": "RdYlBu_r", "Density": "viridis", "Backscatter": "gray_r",
    }
    cmap_name = _CMAP.get(surf_name, "jet")

    if mappable is None:
        norm = plt.Normalize(vmin=vmin, vmax=vmax)
        sm = cm.ScalarMappable(cmap=cmap_name, norm=norm)
        sm.set_array([])
        cbar = plt.colorbar(sm, cax=cax)
    else:
        cbar = plt.colorbar(mappable, cax=cax)

    label = band_labels.get(surf_name, surf_name)
    cbar.set_label(label, fontsize=22, color="#1E3A5F",
                   fontweight="bold", labelpad=18, rotation=270, va="bottom")

    # Smart tick formatting
    val_range = abs(vmax - vmin)
    if val_range > 100:
        fmt = mticker.FuncFormatter(lambda x, _: f"{x:,.0f}")
    elif val_range > 1:
        fmt = mticker.FuncFormatter(lambda x, _: f"{x:.1f}")
    else:
        fmt = mticker.FuncFormatter(lambda x, _: f"{x:.3f}")
    cbar.ax.yaxis.set_major_formatter(fmt)

    cbar.ax.tick_params(labelsize=18, colors="#1A202C",
                        width=1.5, length=5, direction="out", pad=6)
    cbar.outline.set_linewidth(1.2)
    cbar.outline.set_edgecolor("#CBD5E0")


# ── Primitives ─────────────────────────────────────────────────


def _rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, text, left, top, width, height,
          size=14, bold=False, italic=False, color=_TEXT,
          align=PP_ALIGN.LEFT, spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Pretendard"
    p.alignment = align
    if spacing:
        p.font._element.attrib[qn("a:spc")] = str(spacing * 100)


# ── Surface Image Finder ───────────────────────────────────────


def _render_surface_png(surface_dir: Path, surf_name: str) -> Path | None:
    """Render surface image from raw data with matching colorbar.

    Strategy (priority order):
      0. Return existing PNG (created by _render_all_from_gsf) — preferred
      1. Single-band float TIF → render with colorbar
      2. RGBA TIF → display as-is
      3. Fall back to existing JPG
    """
    name_lower = surf_name.lower()

    # ── Strategy 0: Existing PNG already rendered from GSF ──
    for pattern in [surf_name, f"surface_{name_lower}", name_lower]:
        p = surface_dir / f"{pattern}.png"
        if p.exists():
            return p

    _BAND_LABEL = {"Depth": "Depth (m)", "Std_Dev": "Std Dev (m)",
                   "TVU": "TVU (m)", "THU": "THU (m)",
                   "Density": "Soundings / cell",
                   "Backscatter": "Backscatter (dB)"}

    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

            # ── Strategy 1: Single-band float TIF (exact colorbar match) ──
        # Build candidate list with aliases
        _ALIASES = {
            "TVU": ["tvu", "uncertainty"],
            "THU": ["thu"],
            "Depth": ["depth", "dtm", "bat"],
        }
        search_names = [name_lower] + _ALIASES.get(surf_name, [])

        candidates = []
        for sn in search_names:
            candidates += [
                surface_dir / f"{surf_name}.tif",
                surface_dir / f"{surf_name}.tiff",
                surface_dir / f"surface_{sn}.tif",
                surface_dir / f"surface_{sn}.tiff",
            ]
        if surf_name == "Depth":
            candidates += [surface_dir / "DTM.tif", surface_dir / "DTM.tiff",
                            surface_dir / "EDF_BAT_1M.tiff"]
        # Fuzzy match
        for tif in surface_dir.glob("*.tif*"):
            for sn in search_names:
                if sn in tif.stem.lower() and tif not in candidates:
                    candidates.append(tif)

        for tiff_path in candidates:
            if not tiff_path.exists():
                continue
            try:
                import rasterio
                with rasterio.open(str(tiff_path)) as src:
                    if src.count == 1:
                        # Single-band float → render directly (colorbar matches)
                        data = src.read(1).astype(float)
                        nd = src.nodata
                        if nd is not None:
                            data[data == nd] = np.nan
                        data[np.abs(data) > 1e6] = np.nan
                        if np.all(np.isnan(data)):
                            continue
                        return _render_array(plt, np, data, surf_name,
                                             surface_dir, _BAND_LABEL)

                    elif src.count >= 3:
                        # RGBA image — display as-is WITHOUT colorbar
                        return _render_rgba_image(plt, np, tiff_path,
                                                  surface_dir, surf_name)
            except ImportError:
                return None

    except Exception:
        pass

    # ── Strategy 2: Reconstruct from ASCII bands + depth TIF mask ──
    # ── Strategy 2: Reconstruct from ASCII bands ──
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        # Density from Depth grid ASCII bands
        bands_txt = surface_dir / "all_bands.txt"
        raw_tif = surface_dir / "surface_depth.tif"
        if not raw_tif.exists():
            raw_tif = surface_dir / "raw_values.tif"
        if bands_txt.exists() and raw_tif.exists() and surf_name == "Density":
            result = _render_from_ascii_bands(
                plt, np, surface_dir, surf_name, bands_txt, raw_tif, _BAND_LABEL)
            if result:
                return result

        # TVU/THU from TPU grid ASCII bands
        tpu_txt = surface_dir / "tpu_bands.txt"
        tpu_tif = surface_dir / "surface_tpu.tif"
        if tpu_txt.exists() and tpu_tif.exists() and surf_name in ("TVU", "THU"):
            _TPU_COLS = {"TVU": 1, "THU": 2}  # col0=Depth, col1=Depth_TPU, col2=Position_TPU
            result = _render_from_ascii_bands(
                plt, np, surface_dir, surf_name, tpu_txt, tpu_tif, _BAND_LABEL,
                col_override=_TPU_COLS.get(surf_name))
            if result:
                return result
    except Exception:
        pass

    # ── Fallback: Existing JPG ──
    for pattern in [surf_name, f"surface_{name_lower}", name_lower]:
        for ext in [".jpg", ".jpeg"]:
            p = surface_dir / f"{pattern}{ext}"
            if p.exists():
                return p

    return None


def _render_from_ascii_bands(plt, np, surface_dir, surf_name, bands_txt, raw_tif, labels,
                             col_override=None):
    """Reconstruct per-band 2D array from ExportCoverageToASCII output."""
    import rasterio

    _BAND_COLS = {"Depth": 0, "Density": 1, "Std_Dev": 2}
    col_idx = col_override if col_override is not None else _BAND_COLS.get(surf_name)
    if col_idx is None:
        return None

    try:
        ascii_data = np.loadtxt(str(bands_txt))
        if col_idx >= ascii_data.shape[1]:
            return None

        with rasterio.open(str(raw_tif)) as src:
            template = src.read(1).astype(float)
            nd = src.nodata
            if nd is not None:
                template[template == nd] = np.nan
            template[np.abs(template) > 1e6] = np.nan

        # Reconstruct 2D array: fill valid pixels with band values
        band_values = ascii_data[:, col_idx]
        result = np.full_like(template, np.nan)
        valid_mask = ~np.isnan(template)

        if valid_mask.sum() != len(band_values):
            return None

        result[valid_mask] = band_values
        return _render_array(plt, np, result, surf_name, surface_dir, labels)
    except Exception:
        return None


def _render_array(plt, np, data, surf_name, surface_dir, labels):
    """Render a 2D numpy array with colorbar — image and colorbar use same colormap."""
    _CMAP = {
        "Depth": "jet", "Std_Dev": "YlOrRd", "TVU": "RdYlBu_r",
        "THU": "RdYlBu_r", "Density": "viridis", "Backscatter": "gray_r",
    }

    # CARIS exports depth as negative (below surface) — flip for display
    if surf_name == "Depth":
        vv = data[~np.isnan(data)]
        if len(vv) > 0 and vv.mean() < 0:
            data = np.abs(data)

    # Tight crop
    valid = ~np.isnan(data)
    rows_valid = np.any(valid, axis=1)
    cols_valid = np.any(valid, axis=0)
    if rows_valid.any() and cols_valid.any():
        r0, r1 = np.where(rows_valid)[0][[0, -1]]
        c0, c1 = np.where(cols_valid)[0][[0, -1]]
        pad = max(2, int(0.02 * max(r1 - r0, c1 - c0)))
        r0, r1 = max(0, r0 - pad), min(data.shape[0], r1 + pad)
        c0, c1 = max(0, c0 - pad), min(data.shape[1], c1 + pad)
        data = data[r0:r1, c0:c1]

    # Rotate if very elongated vertically (for landscape PPT)
    nrows, ncols = data.shape
    if nrows > ncols * 2.5:
        data = np.rot90(data, k=1)

    # Percentile clipping for colorbar
    valid_vals = data[~np.isnan(data)]
    vmin = float(np.percentile(valid_vals, 2))
    vmax = float(np.percentile(valid_vals, 98))
    if abs(vmax - vmin) < 1e-10:
        vmin, vmax = float(valid_vals.min()), float(valid_vals.max())

    # Adaptive figsize
    nrows, ncols = data.shape
    aspect = nrows / max(ncols, 1)
    fig_h = 8
    fig_w = max(10, fig_h / max(aspect, 0.3))
    fig_w = min(fig_w, 20) + 2

    cmap_obj = plt.get_cmap(_CMAP.get(surf_name, "jet")).copy()
    cmap_obj.set_bad(color="white")

    fig, (ax, cax) = plt.subplots(1, 2, figsize=(fig_w, fig_h),
        gridspec_kw={"width_ratios": [16, 1], "wspace": 0.04},
        facecolor="white")
    im = ax.imshow(data, cmap=cmap_obj, aspect="equal", interpolation="nearest",
                   vmin=vmin, vmax=vmax)
    _draw_colorbar(plt, cax, vmin, vmax,
                   surf_name, labels, mappable=im)
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)

    png_path = surface_dir / f"{surf_name}.png"
    fig.savefig(str(png_path), dpi=200, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    return png_path


def _render_rgba_image(plt, np, tiff_path, surface_dir, surf_name):
    """Display RGBA GeoTIFF with colorbar from band_stats.json."""
    import rasterio

    _BAND_LABEL = {"Depth": "Depth (m)", "Std_Dev": "Std Dev (m)",
                   "TVU": "TVU (m)", "THU": "THU (m)",
                   "Density": "Density (pts/cell)",
                   "Backscatter": "Backscatter (dB)"}

    with rasterio.open(str(tiff_path)) as src:
        if src.count == 4:
            rgba = src.read([1, 2, 3, 4])
            rgba = np.moveaxis(rgba, 0, -1).astype(np.uint8)
            black = (rgba[:,:,0] == 0) & (rgba[:,:,1] == 0) & (rgba[:,:,2] == 0)
            rgba[black] = [255, 255, 255, 255]
            img_arr = rgba[:,:,:3]
        else:
            rgb = src.read([1, 2, 3])
            img_arr = np.moveaxis(rgb, 0, -1).astype(np.uint8)
            black = (img_arr[:,:,0] == 0) & (img_arr[:,:,1] == 0) & (img_arr[:,:,2] == 0)
            img_arr[black] = [255, 255, 255]

    # Crop to content
    non_white = np.any(img_arr < 250, axis=2)
    if non_white.any():
        rv = np.any(non_white, axis=1)
        cv = np.any(non_white, axis=0)
        r0, r1 = np.where(rv)[0][[0, -1]]
        c0, c1 = np.where(cv)[0][[0, -1]]
        pad = max(4, int(0.02 * max(r1 - r0, c1 - c0)))
        r0 = max(0, r0 - pad); r1 = min(img_arr.shape[0], r1 + pad)
        c0 = max(0, c0 - pad); c1 = min(img_arr.shape[1], c1 + pad)
        img_arr = img_arr[r0:r1, c0:c1]

    # Try to get value range for colorbar from band_stats.json
    vmin, vmax = None, None
    stats_file = surface_dir / "band_stats.json"
    if stats_file.is_file():
        try:
            import json
            with open(str(stats_file)) as f:
                stats = json.load(f)
            if surf_name in stats:
                vmin = stats[surf_name]["min"]
                vmax = stats[surf_name]["max"]
        except Exception:
            pass

    has_cbar = vmin is not None and vmax is not None and abs(vmax - vmin) > 1e-10

    if has_cbar:
        fig, (ax, cax) = plt.subplots(1, 2, figsize=(14, 10),
            gridspec_kw={"width_ratios": [16, 1], "wspace": 0.04},
            facecolor="white")
    else:
        fig, ax = plt.subplots(figsize=(14, 10), facecolor="white")
        cax = None

    ax.imshow(img_arr, aspect="equal")
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)

    if has_cbar and cax is not None:
        _draw_colorbar(plt, cax, vmin, vmax, surf_name, _BAND_LABEL)

    png_path = surface_dir / f"{surf_name}.png"
    fig.savefig(str(png_path), dpi=200, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    return png_path


def _render_track_plot(gsf_main, output_dir: Path | None = None) -> tuple | None:
    """Generate track plot PNG from GSF ping lat/lon coordinates."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker
        import numpy as np

        lats = [p.latitude for p in gsf_main.pings if p.latitude != 0]
        lons = [p.longitude for p in gsf_main.pings if p.longitude != 0]

        if len(lats) < 2:
            return None

        lats = np.array(lats)
        lons = np.array(lons)

        # Convert WGS84 lat/lon to UTM
        plot_x, plot_y = lons, lats
        utm_label = ""
        try:
            from pyproj import Transformer
            mean_lon, mean_lat = lons.mean(), lats.mean()
            zone = int((mean_lon + 180) / 6) + 1
            hemisphere = "north" if mean_lat >= 0 else "south"
            epsg = 32600 + zone if hemisphere == "north" else 32700 + zone
            if 124 < mean_lon < 132 and 33 < mean_lat < 39:
                epsg = 32652
            transformer = Transformer.from_crs("EPSG:4326", f"EPSG:{epsg}", always_xy=True)
            plot_x, plot_y = transformer.transform(lons, lats)
            utm_label = f"EPSG:{epsg}"
        except Exception:
            pass

        is_utm = abs(plot_x.mean()) > 1000

        # Adaptive figure size — match data aspect ratio to fill slide
        xr = plot_x.max() - plot_x.min()
        yr = plot_y.max() - plot_y.min()
        data_aspect = xr / yr if yr > 0 else 1.5
        # Target slide area is ~11.7" x 5.3" (ratio ~2.2)
        if data_aspect > 1.5:
            fig_w, fig_h = 14, 8
        elif data_aspect > 0.5:
            fig_w, fig_h = 12, 10
        else:
            fig_w, fig_h = 9, 12

        fig, ax = plt.subplots(figsize=(fig_w, fig_h), facecolor="white")

        # Gradient track line (segments colored by time)
        colors = np.linspace(0, 1, len(plot_x))

        # Thicker track backbone
        ax.plot(plot_x, plot_y, "-", color="#CBD5E0", linewidth=3, alpha=0.4, zorder=1)

        # Color scatter — larger points for visibility
        sc = ax.scatter(plot_x, plot_y, c=colors, cmap="RdYlBu_r",
                        s=8, alpha=0.85, edgecolors="none", zorder=3)

        # Start/End markers — larger, with labels
        ax.plot(plot_x[0], plot_y[0], "o", color="#38A169", markersize=16,
                markeredgecolor="white", markeredgewidth=2.5,
                label="Start", zorder=5)
        ax.plot(plot_x[-1], plot_y[-1], "s", color="#E53E3E", markersize=16,
                markeredgecolor="white", markeredgewidth=2.5,
                label="End", zorder=5)

        if is_utm:
            ax.set_xlabel("Easting (m)", fontsize=15, color="#1A202C",
                          fontweight="semibold", labelpad=12)
            ax.set_ylabel("Northing (m)", fontsize=15, color="#1A202C",
                          fontweight="semibold", labelpad=12)
            ax.xaxis.set_major_formatter(mticker.FuncFormatter(
                lambda x, _: f"{x:,.0f}"))
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(
                lambda y, _: f"{y:,.0f}"))
            if utm_label:
                ax.text(0.98, 0.98, utm_label, transform=ax.transAxes,
                        fontsize=11, color="#A0AEC0", ha="right", va="top",
                        fontstyle="italic")
            pad_m = max(xr, yr) * 0.15 + 30
            ax.set_xlim(plot_x.min() - pad_m, plot_x.max() + pad_m)
            ax.set_ylim(plot_y.min() - pad_m, plot_y.max() + pad_m)
        else:
            ax.xaxis.set_major_formatter(mticker.FuncFormatter(
                lambda x, _: f"{x:.4f}\u00b0"))
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(
                lambda y, _: f"{y:.4f}\u00b0"))
            ax.set_xlabel("Longitude", fontsize=15, color="#1A202C",
                          fontweight="semibold", labelpad=12)
            ax.set_ylabel("Latitude", fontsize=15, color="#1A202C",
                          fontweight="semibold", labelpad=12)
            lat_range = lats.max() - lats.min()
            lon_range = lons.max() - lons.min()
            pad_d = max(lat_range, lon_range) * 0.10
            ax.set_xlim(lons.min() - pad_d, lons.max() + pad_d)
            ax.set_ylim(lats.min() - pad_d, lats.max() + pad_d)

        ax.tick_params(labelsize=13, colors="#4A5568", width=1.2, length=5)
        plt.setp(ax.get_xticklabels(), rotation=30, ha="right")
        ax.set_aspect("equal")

        # Legend — styled box
        legend = ax.legend(fontsize=13, loc="upper right", framealpha=0.95,
                           edgecolor="#CBD5E0", fancybox=True, shadow=False,
                           handletextpad=0.8, borderpad=0.8)
        legend.get_frame().set_linewidth(1.2)

        # Grid
        ax.grid(True, alpha=0.2, color="#A0AEC0", linewidth=0.6, linestyle="--")
        for spine in ax.spines.values():
            spine.set_color("#CBD5E0")
            spine.set_linewidth(1.2)

        # Total km badge
        total_km = _compute_line_km(lats, lons)
        ax.text(0.02, 0.03, f"Total: {total_km:.1f} km",
                transform=ax.transAxes, fontsize=13, color="#1E3A5F",
                fontweight="bold", ha="left", va="bottom",
                bbox=dict(boxstyle="round,pad=0.4", facecolor="#F0F4F8",
                          edgecolor="#CBD5E0", linewidth=1.2, alpha=0.95))

        save_dir = output_dir or Path(".")
        save_dir.mkdir(parents=True, exist_ok=True)
        png_path = save_dir / "trackplot.png"
        fig.savefig(str(png_path), dpi=200, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        plt.close(fig)
        return png_path, total_km

    except Exception:
        return None


def _compute_line_km(lats, lons):
    """Compute total track line distance in km using Haversine formula."""
    import numpy as np
    R = 6371.0
    lat_r = np.radians(lats)
    lon_r = np.radians(lons)
    dlat = np.diff(lat_r)
    dlon = np.diff(lon_r)
    a = np.sin(dlat / 2) ** 2 + np.cos(lat_r[:-1]) * np.cos(lat_r[1:]) * np.sin(dlon / 2) ** 2
    c = 2 * np.arctan2(np.sqrt(a), np.sqrt(1 - a))
    return float(np.sum(R * c))
