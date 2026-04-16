"""MBESQC ExportService -- Excel/Word/PPT report generation worker."""

from __future__ import annotations

import json
import os
import tempfile
from datetime import datetime
from pathlib import Path

from PySide6.QtCore import QObject, Signal, Slot

from desktop.services.data_service import DataService
from desktop.services.insight_service import (
    build_action_checklist,
    build_history_story,
    build_issue_spotlight,
    build_module_rows,
    build_project_context,
    build_result_overview,
    build_run_diff,
    build_settings_assistant,
    describe_result_snapshot,
    format_number,
)


def _extract_export_provenance(latest_data: dict | None) -> dict:
    """Pull the persisted provenance payload back out for report surfaces."""
    data = latest_data if isinstance(latest_data, dict) else {}
    summary = data.get("provenance_summary")
    if not isinstance(summary, dict):
        summary = DataService.extract_provenance_summary(data)

    manifest = data.get("provenance_manifest")
    if not isinstance(manifest, dict):
        manifest = DataService.extract_provenance_manifest(data)

    raw = {}

    offset_validation = data.get("offset_validation") or {}
    if isinstance(offset_validation, dict):
        raw = offset_validation.get("provenance") or {}

    if not raw:
        raw = data.get("provenance") or {}

    if not isinstance(raw, dict):
        raw = {}

    om = raw.get("om") or {}
    if not isinstance(om, dict):
        om = {}

    semantic = om.get("semantic") or {}
    if not isinstance(semantic, dict):
        semantic = {}

    semantic_checks = semantic.get("checks") or []
    if not isinstance(semantic_checks, list):
        semantic_checks = []

    decision = om.get("decision") or {}
    if not isinstance(decision, dict):
        decision = {}

    resolution_chain = om.get("resolution_chain") or raw.get("resolution_chain") or []
    if not isinstance(resolution_chain, list):
        resolution_chain = []

    return _redact_export_provenance_paths({
        "raw": raw,
        "summary": summary,
        "om": om,
        "decision": decision,
        "semantic": semantic,
        "semantic_checks": semantic_checks,
        "resolution_chain": resolution_chain,
        "manifest": manifest,
        "has_data": bool(summary.get("has_data") or raw or semantic_checks),
    })


def _redact_export_path(value) -> str:
    """Collapse a local filesystem path to its basename for shared exports."""
    if value in (None, ""):
        return ""
    if isinstance(value, dict):
        return _redact_export_provenance_paths(value)
    if isinstance(value, list):
        return [_redact_export_path(item) for item in value]
    try:
        return Path(str(value)).name or str(value)
    except Exception:
        return str(value)


def _redact_export_uri(value) -> str:
    """Collapse a URI-like value to a non-sensitive placeholder."""
    if value in (None, ""):
        return ""
    if isinstance(value, dict):
        return _redact_export_provenance_paths(value)
    if isinstance(value, list):
        return [_redact_export_uri(item) for item in value]
    text = str(value)
    if "://" in text:
        return "<redacted:uri>"
    return _redact_export_path(text)


def _is_export_path_key(key) -> bool:
    key_text = str(key).strip().lower()
    collapsed = key_text.replace("_", "")
    return (
        key_text in {"path", "filepath", "file_path", "db_path", "source_path", "source_filepath"}
        or collapsed in {"path", "filepath", "dbpath", "sourcepath", "sourcefilepath"}
        or key_text.endswith("_path")
        or key_text.endswith("_filepath")
    )


def _is_export_uri_key(key) -> bool:
    key_text = str(key).strip().lower()
    collapsed = key_text.replace("_", "")
    return (
        key_text in {"uri", "url", "source_uri", "source_url"}
        or collapsed in {"uri", "url", "sourceuri", "sourceurl"}
        or key_text.endswith("_uri")
        or key_text.endswith("_url")
    )


def _redact_export_provenance_paths(value):
    """Redact path-bearing fields from a provenance payload copy."""
    if isinstance(value, dict):
        redacted = {}
        for key, item in value.items():
            if _is_export_path_key(key):
                redacted[key] = _redact_export_path(item)
            elif _is_export_uri_key(key):
                redacted[key] = _redact_export_uri(item)
            else:
                redacted[key] = _redact_export_provenance_paths(item)
        return redacted
    if isinstance(value, list):
        return [_redact_export_provenance_paths(item) for item in value]
    return value


def _format_provenance_summary_value(summary: dict) -> str:
    source = summary.get("source") or "unknown"
    state = summary.get("semantic_state") or "unspecified"
    checks_total = int(summary.get("semantic_checks_total") or 0)
    checks_passed = int(summary.get("semantic_checks_passed") or 0)
    bits = [source, state]
    if checks_total:
        bits.append(f"checks {checks_passed}/{checks_total}")
    fallback_scope = summary.get("fallback_scope") or ""
    if fallback_scope:
        bits.append(f"scope {fallback_scope}")
    return " / ".join(bits)


def _format_export_provenance_lines(provenance: dict) -> list[str]:
    lines: list[str] = []
    summary = provenance.get("summary") or {}
    decision = provenance.get("decision") or {}
    semantic = provenance.get("semantic") or {}
    semantic_checks = provenance.get("semantic_checks") or []
    resolution_chain = provenance.get("resolution_chain") or []

    if summary.get("has_data"):
        lines.append("Persisted provenance summary: " + _format_provenance_summary_value(summary))

    source = decision.get("source") or "unknown"
    path = decision.get("path") or ""
    if source != "unknown" or path:
        decision_bits = [f"OffsetManager source: {source}"]
        if path:
            decision_bits.append(f"path {path}")
        lines.append(" | ".join(decision_bits))

    state = semantic.get("state") or ""
    hint = semantic.get("hint") or ""
    if state:
        lines.append(f"Semantic state: {state}")
    if hint:
        lines.append(f"Semantic hint: {hint}")

    if semantic_checks:
        passed = sum(1 for check in semantic_checks if check.get("passed"))
        lines.append(f"Semantic checks: {passed}/{len(semantic_checks)} passed")
        for check in semantic_checks[:4]:
            status = "PASS" if check.get("passed") else "FAIL"
            name = check.get("name") or "check"
            detail = check.get("detail") or ""
            line = f"{name} [{status}]"
            if detail:
                line += f": {detail}"
            lines.append(line)

    if resolution_chain:
        lines.append(f"Resolution chain steps: {len(resolution_chain)}")

    if not lines:
        lines.append("No persisted provenance payload was available in the exported snapshot.")

    return lines


class ExportWorker(QObject):
    """Background worker for report generation."""

    finished = Signal(str)   # output path
    error = Signal(str)      # error message

    def __init__(self, project_id: int, fmt: str, output_path: str):
        super().__init__()
        self._project_id = project_id
        self._fmt = fmt
        self._output_path = output_path

    @Slot()
    def run(self):
        try:
            if self._fmt == "excel":
                self._export_excel()
            elif self._fmt == "word":
                self._export_word()
            elif self._fmt == "ppt":
                self._export_ppt()
            else:
                self.error.emit(f"Unsupported format: {self._fmt}")
                return

            DataService.log_activity(
                self._project_id,
                "export_complete",
                f"{self._fmt}:{Path(self._output_path).name}",
            )
            self.finished.emit(self._output_path)

        except Exception as e:
            self.error.emit(str(e))

    def _load_export_context(self):
        project = DataService.get_project(self._project_id)
        results = DataService.get_project_qc_results(self._project_id)
        done_results = [r for r in results if r.get("status") == "done"]
        if not done_results:
            raise ValueError("완료된 QC 결과가 없습니다. 먼저 프로젝트 QC를 실행하세요.")

        latest = max(done_results, key=lambda r: r.get("finished_at") or "")
        try:
            latest_data = latest.get("result_payload") or json.loads(latest.get("result_json") or "{}")
        except (TypeError, json.JSONDecodeError):
            latest_data = {}

        files = DataService.get_project_files(self._project_id)
        file_counts = {
            "pds_count": sum(1 for f in files if f.get("format", "").lower() == "pds"),
            "gsf_count": sum(1 for f in files if f.get("format", "").lower() == "gsf"),
            "hvf_count": 1 if project and project.get("hvf_dir") else 0,
        }
        context = build_project_context(project, latest_result=latest, file_counts=file_counts)
        overview = build_result_overview(latest_data)
        module_rows = build_module_rows(latest_data)
        issue_spotlight = build_issue_spotlight(latest_data)
        action_items = build_action_checklist(latest_data)
        history_story = build_history_story(done_results)
        run_diff = build_run_diff(done_results)
        settings = build_settings_assistant(project, latest_data, file_counts)
        provenance = _extract_export_provenance(latest_data)
        history_rows = []
        for result in done_results:
            anchor_file = DataService.get_file(result.get("file_id", 0)) if result.get("file_id") else None
            history_rows.append({
                "snapshot": describe_result_snapshot(result, anchor_file=anchor_file),
                "score": format_number(result.get("score", 0), 1),
                "grade": result.get("grade", ""),
                "status": result.get("status", ""),
                "completed": (result.get("finished_at", "") or "")[:19],
            })
        return (
            project, results, done_results, latest, latest_data,
            context, overview, module_rows, history_rows, issue_spotlight, action_items, history_story,
            run_diff, settings, provenance,
        )

    def _export_excel(self):
        """Generate Excel report with charts embedded."""
        import openpyxl
        from openpyxl.styles import Alignment, Font as XlFont, PatternFill, Border, Side
        from openpyxl.drawing.image import Image as XlImage

        (
            project, results, done_results, latest, data,
            context, overview, module_rows, history_rows, issue_spotlight, action_items, history_story,
            run_diff, settings, provenance,
        ) = self._load_export_context()

        wb = openpyxl.Workbook()
        ws_exec = wb.active
        ws_exec.title = "Executive Summary"

        # Styles
        h_fill = PatternFill(start_color="1A2236", end_color="1A2236", fill_type="solid")
        h_font = XlFont(name="Pretendard", size=11, bold=True, color="FFFFFF")
        body_font = XlFont(name="Pretendard", size=11)
        title_font = XlFont(name="Pretendard", size=16, bold=True, color="10B981")
        wrap = Alignment(wrap_text=True, vertical="top")
        thin = Border(
            left=Side(style="thin"), right=Side(style="thin"),
            top=Side(style="thin"), bottom=Side(style="thin"),
        )

        # Executive summary sheet
        ws_exec.merge_cells("A1:E1")
        ws_exec["A1"] = f"MBES QC Report -- {project['name']}" if project else "MBES QC Report"
        ws_exec["A1"].font = title_font
        ws_exec["A2"] = f"Vessel: {project.get('vessel', '')}" if project else ""
        ws_exec["A2"].font = XlFont(name="Pretendard", size=11, color="6B7280")

        exec_rows = [
            ("Current focus", overview["headline"]),
            ("Latest snapshot", context["snapshot_text"]),
            ("QC overview", overview["body"]),
            ("Trend", history_story["headline"]),
            ("Trend detail", history_story["body"]),
            ("Run-to-run diff", run_diff["body"]),
            ("Critical findings", "\n".join(
                f"[{item['status']}] {item['module']} - {item['title']}: {item['evidence']}"
                for item in issue_spotlight[:4]
            ) or "---"),
            ("Recommended actions", "\n".join(f"- {line}" for line in action_items[:4]) or "---"),
            ("Settings", " / ".join(settings.get("current_state", [])[:4]) or "---"),
            ("Settings guide", "\n".join(f"- {line}" for line in settings.get("recommendations", [])[:4]) or "---"),
            ("Project flow", context["flow"]),
            ("Offset reference", context["offset_text"]),
            ("Readiness", context["readiness_text"]),
            ("Next checks", "\n".join(f"- {line}" for line in overview.get("next_steps", [])[:4])),
        ]
        for idx, (label, value) in enumerate(exec_rows, 4):
            ws_exec.cell(row=idx, column=1, value=label).font = h_font
            ws_exec.cell(row=idx, column=1).fill = h_fill
            ws_exec.cell(row=idx, column=1).border = thin
            ws_exec.cell(row=idx, column=2, value=value).font = body_font
            ws_exec.cell(row=idx, column=2).alignment = wrap
            ws_exec.cell(row=idx, column=2).border = thin
            ws_exec.row_dimensions[idx].height = 48
        ws_exec.column_dimensions["A"].width = 20
        ws_exec.column_dimensions["B"].width = 95

        # Results table
        ws = wb.create_sheet("QC Summary")
        ws.merge_cells("A1:F1")
        ws["A1"] = "QC Run History"
        ws["A1"].font = title_font
        ws["A2"] = "이 export는 최신 완료 프로젝트 QC 스냅샷의 해설을 기준으로 작성되었습니다."
        ws["A2"].font = XlFont(name="Pretendard", size=10, color="6B7280")

        row = 4
        headers = ["Snapshot", "Score", "Grade", "Status", "Completed"]
        for col, h in enumerate(headers, 1):
            cell = ws.cell(row=row, column=col, value=h)
            cell.font = h_font
            cell.fill = h_fill
            cell.border = thin

        for i, history in enumerate(history_rows, row + 1):
            ws.cell(row=i, column=1, value=history["snapshot"]).font = body_font
            ws.cell(row=i, column=2, value=history["score"]).font = body_font
            ws.cell(row=i, column=3, value=history["grade"]).font = body_font
            ws.cell(row=i, column=4, value=history["status"]).font = body_font
            ws.cell(row=i, column=5, value=history["completed"]).font = body_font

            for c in range(1, 6):
                ws.cell(row=i, column=c).border = thin

        ws.column_dimensions["A"].width = 60
        ws.column_dimensions["B"].width = 12
        ws.column_dimensions["C"].width = 10
        ws.column_dimensions["D"].width = 12
        ws.column_dimensions["E"].width = 22

        ws_modules = wb.create_sheet("Module Guidance")
        module_headers = ["Module", "Verdict", "Why it matters", "Current reading", "Next check"]
        for col, h in enumerate(module_headers, 1):
            cell = ws_modules.cell(row=1, column=col, value=h)
            cell.font = h_font
            cell.fill = h_fill
            cell.border = thin
        for row_idx, story in enumerate(module_rows, 2):
            values = [
                story["module"],
                story["status"],
                story["importance"],
                story["current_reading"],
                " / ".join(story.get("next_steps", [])[:3]),
            ]
            for col_idx, value in enumerate(values, 1):
                cell = ws_modules.cell(row=row_idx, column=col_idx, value=value)
                cell.font = body_font
                cell.alignment = wrap
                cell.border = thin
            ws_modules.row_dimensions[row_idx].height = 56
        for col, width in {"A": 18, "B": 10, "C": 34, "D": 48, "E": 38}.items():
            ws_modules.column_dimensions[col].width = width

        if issue_spotlight:
            ws_findings = wb.create_sheet("Issue Spotlight")
            finding_headers = ["Module", "Verdict", "Finding", "Evidence", "Recommended action"]
            for col, header in enumerate(finding_headers, 1):
                cell = ws_findings.cell(row=1, column=col, value=header)
                cell.font = h_font
                cell.fill = h_fill
                cell.border = thin
            for row_idx, finding in enumerate(issue_spotlight, 2):
                values = [
                    finding["module"],
                    finding["status"],
                    finding["title"],
                    finding["evidence"],
                    finding["action"],
                ]
                for col_idx, value in enumerate(values, 1):
                    cell = ws_findings.cell(row=row_idx, column=col_idx, value=value)
                    cell.font = body_font
                    cell.alignment = wrap
                    cell.border = thin
                ws_findings.row_dimensions[row_idx].height = 54
            for col, width in {"A": 18, "B": 10, "C": 28, "D": 48, "E": 42}.items():
                ws_findings.column_dimensions[col].width = width

        if provenance.get("has_data"):
            ws_prov = wb.create_sheet("Provenance")
            ws_prov.merge_cells("A1:E1")
            ws_prov["A1"] = "QC Provenance"
            ws_prov["A1"].font = title_font
            ws_prov["A2"] = "Persisted OffsetManager provenance and semantic checks from the latest completed QC snapshot."
            ws_prov["A2"].font = XlFont(name="Pretendard", size=10, color="6B7280")

            summary = provenance.get("summary") or {}
            prov_rows = [
                ("Persisted summary", _format_provenance_summary_value(summary)) if summary.get("has_data") else None,
                ("Decision source", provenance["decision"].get("source", "")),
                ("Decision path", provenance["decision"].get("path", "")),
                ("Semantic state", provenance["semantic"].get("state", "")),
                ("Semantic hint", provenance["semantic"].get("hint", "")),
                ("Resolution chain steps", str(len(provenance.get("resolution_chain") or []))),
            ]
            for row_idx, entry in enumerate([row for row in prov_rows if row], 4):
                label, value = entry
                ws_prov.cell(row=row_idx, column=1, value=label).font = h_font
                ws_prov.cell(row=row_idx, column=1).fill = h_fill
                ws_prov.cell(row=row_idx, column=1).border = thin
                ws_prov.cell(row=row_idx, column=2, value=value).font = body_font
                ws_prov.cell(row=row_idx, column=2).alignment = wrap
                ws_prov.cell(row=row_idx, column=2).border = thin

            check_row = 10
            ws_prov.cell(row=check_row, column=1, value="Semantic checks").font = h_font
            ws_prov.cell(row=check_row, column=1).fill = h_fill
            ws_prov.cell(row=check_row, column=1).border = thin
            ws_prov.cell(row=check_row, column=2, value="Passed").font = h_font
            ws_prov.cell(row=check_row, column=2).fill = h_fill
            ws_prov.cell(row=check_row, column=2).border = thin
            ws_prov.cell(row=check_row, column=3, value="Blocking").font = h_font
            ws_prov.cell(row=check_row, column=3).fill = h_fill
            ws_prov.cell(row=check_row, column=3).border = thin
            ws_prov.cell(row=check_row, column=4, value="Severity").font = h_font
            ws_prov.cell(row=check_row, column=4).fill = h_fill
            ws_prov.cell(row=check_row, column=4).border = thin
            ws_prov.cell(row=check_row, column=5, value="Detail").font = h_font
            ws_prov.cell(row=check_row, column=5).fill = h_fill
            ws_prov.cell(row=check_row, column=5).border = thin

            for row_idx, check in enumerate(provenance.get("semantic_checks") or [], check_row + 1):
                values = [
                    check.get("name", ""),
                    "YES" if check.get("passed") else "NO",
                    "YES" if check.get("blocking") else "NO",
                    check.get("severity", ""),
                    check.get("detail", ""),
                ]
                for col_idx, value in enumerate(values, 1):
                    cell = ws_prov.cell(row=row_idx, column=col_idx, value=value)
                    cell.font = body_font
                    cell.alignment = wrap
                    cell.border = thin
                ws_prov.row_dimensions[row_idx].height = 42

            json_row = check_row + max(len(provenance.get("semantic_checks") or []), 1) + 3
            ws_prov.cell(row=json_row, column=1, value="Raw provenance JSON").font = h_font
            ws_prov.cell(row=json_row, column=1).fill = h_fill
            ws_prov.cell(row=json_row, column=1).border = thin
            raw_json = json.dumps(
                _redact_export_provenance_paths(provenance.get("raw") or {}),
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
            ws_prov.cell(row=json_row + 1, column=1, value=raw_json)
            ws_prov.cell(row=json_row + 1, column=1).font = XlFont(name="Consolas", size=10)
            ws_prov.cell(row=json_row + 1, column=1).alignment = wrap
            ws_prov.column_dimensions["A"].width = 28
            ws_prov.column_dimensions["B"].width = 20
            ws_prov.column_dimensions["C"].width = 14
            ws_prov.column_dimensions["D"].width = 12
            ws_prov.column_dimensions["E"].width = 64

        # ── Charts sheet ──
        temp_files = []
        try:
            if data:
                from desktop.services.chart_renderer import render_qc_radar
                ws_charts = wb.create_sheet("Charts")
                row_pos = 1

                scores = {}
                for story in module_rows:
                    if story["status"] == "PASS":
                        scores[story["module"]] = 95
                    elif story["status"] == "WARNING":
                        scores[story["module"]] = 60
                    elif story["status"] == "FAIL":
                        scores[story["module"]] = 20

                if scores:
                    png = render_qc_radar(scores, "QC Score Breakdown")
                    tmp = tempfile.NamedTemporaryFile(
                        delete=False, suffix=".png", prefix="mbesqc_radar_")
                    tmp.write(png)
                    tmp.close()
                    temp_files.append(tmp.name)

                    ws_charts.cell(row=row_pos, column=1, value="QC Score Breakdown")
                    ws_charts.cell(row=row_pos, column=1).font = XlFont(
                        name="Pretendard", size=12, bold=True)
                    row_pos += 1

                    img = XlImage(tmp.name)
                    img.width = 500
                    img.height = 500
                    ws_charts.add_image(img, f"A{row_pos}")
                    row_pos += 28

        except Exception as chart_err:
            import logging
            logging.getLogger(__name__).warning("Chart generation skipped: %s", chart_err)

        try:
            wb.save(self._output_path)
        finally:
            for tmp_path in temp_files:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    def _export_word(self):
        """Generate Word report with python-docx."""
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        (
            project, results, done_results, latest, data,
            context, overview, module_rows, history_rows, issue_spotlight, action_items, history_story,
            run_diff, settings, provenance,
        ) = self._load_export_context()

        doc = Document()

        # Default font for body
        style = doc.styles['Normal']
        style.font.name = 'Pretendard'
        style.font.size = Pt(11)

        # Title
        title = doc.add_heading(level=0)
        run = title.add_run(f"MBES QC Report")
        run.font.size = Pt(24)
        run.font.name = 'Pretendard'

        # Project info
        doc.add_paragraph(f"Project: {project['name']}" if project else "")
        doc.add_paragraph(f"Vessel: {project.get('vessel', '')}" if project else "")
        doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        doc.add_paragraph("")

        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(overview["headline"])
        doc.add_paragraph(overview["body"])
        doc.add_paragraph(f"Trend: {history_story['headline']}")
        doc.add_paragraph(history_story["body"])
        doc.add_paragraph(f"Run-to-run diff: {run_diff['body']}")
        if settings.get("current_state"):
            doc.add_paragraph("Settings: " + " / ".join(settings["current_state"][:4]))
        doc.add_paragraph(f"Project flow: {context['flow']}")
        doc.add_paragraph(f"Latest snapshot: {context['snapshot_text']}")
        doc.add_paragraph(f"Offset reference: {context['offset_text']}")
        doc.add_paragraph("Readiness")
        for line in context["readiness_text"].splitlines():
            doc.add_paragraph(line, style="List Bullet")
        if overview.get("next_steps"):
            doc.add_paragraph("Next checks")
            for line in overview["next_steps"][:4]:
                doc.add_paragraph(line, style="List Bullet")
        if issue_spotlight:
            doc.add_heading("Critical Findings", level=1)
            for finding in issue_spotlight[:5]:
                doc.add_paragraph(
                    f"[{finding['status']}] {finding['module']} - {finding['title']}: {finding['evidence']}",
                    style="List Bullet",
                )
        if action_items:
            doc.add_heading("Recommended Actions", level=1)
            for line in action_items[:5]:
                doc.add_paragraph(line, style="List Bullet")
        if settings.get("recommendations"):
            doc.add_heading("Settings Guidance", level=1)
            for line in settings["recommendations"][:5]:
                doc.add_paragraph(line, style="List Bullet")
        if history_story.get("persistent_modules") or history_story.get("worsened_modules") or history_story.get("improved_modules"):
            doc.add_heading("Trend and Persistence", level=1)
            if history_story.get("persistent_modules"):
                doc.add_paragraph(
                    "Persistent issues: " + ", ".join(history_story["persistent_modules"][:5]),
                    style="List Bullet",
                )
            if history_story.get("worsened_modules"):
                doc.add_paragraph(
                    "Worsened vs previous: " + ", ".join(history_story["worsened_modules"][:5]),
                    style="List Bullet",
                )
            if history_story.get("improved_modules"):
                doc.add_paragraph(
                    "Improved vs previous: " + ", ".join(history_story["improved_modules"][:5]),
                    style="List Bullet",
                )
        if run_diff.get("changes"):
            doc.add_heading("Run-to-Run Module Diff", level=1)
            for item in run_diff["changes"][:5]:
                doc.add_paragraph(
                    f"{item['module']}: {item['previous_status']} -> {item['latest_status']} | {item['latest_reading']}",
                    style="List Bullet",
                )

        if provenance.get("has_data"):
            doc.add_heading("Provenance", level=1)
            for line in _format_export_provenance_lines(provenance):
                doc.add_paragraph(line, style="List Bullet")
            raw_json = json.dumps(
                _redact_export_provenance_paths(provenance.get("raw") or {}),
                ensure_ascii=False,
                indent=2,
                sort_keys=True,
            )
            doc.add_paragraph("Raw provenance JSON")
            doc.add_paragraph(raw_json)

        # Summary table
        if done_results:
            doc.add_heading("QC Results Summary", level=1)
            table = doc.add_table(rows=1 + len(history_rows), cols=5)
            table.style = "Table Grid"

            headers = ["Snapshot", "Score", "Grade", "Status", "Completed"]
            for i, h in enumerate(headers):
                table.cell(0, i).text = h

            for i, history in enumerate(history_rows):
                table.cell(i + 1, 0).text = history["snapshot"]
                table.cell(i + 1, 1).text = history["score"]
                table.cell(i + 1, 2).text = history["grade"]
                table.cell(i + 1, 3).text = history["status"]
                table.cell(i + 1, 4).text = history["completed"]

        if module_rows:
            doc.add_heading("Module Guidance", level=1)
            table = doc.add_table(rows=1 + len(module_rows), cols=5)
            table.style = "Table Grid"
            headers = ["Module", "Verdict", "Why it matters", "Current reading", "Next check"]
            for i, header in enumerate(headers):
                table.cell(0, i).text = header
            for row_idx, story in enumerate(module_rows, 1):
                table.cell(row_idx, 0).text = story["module"]
                table.cell(row_idx, 1).text = story["status"]
                table.cell(row_idx, 2).text = story["importance"]
                table.cell(row_idx, 3).text = story["current_reading"]
                table.cell(row_idx, 4).text = " / ".join(story.get("next_steps", [])[:3])

        doc.save(self._output_path)

    def _export_ppt(self):
        """Generate PPT DQR with python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        (
            project, results, done_results, latest, data,
            context, overview, module_rows, history_rows, issue_spotlight, action_items, history_story,
            run_diff, settings, provenance,
        ) = self._load_export_context()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "MBES Data Quality Report"
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
        p.font.bold = True

        p2 = tf.add_paragraph()
        p2.text = f"{project['name']} | {project.get('vessel', '')}" if project else ""
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        p3 = tf.add_paragraph()
        p3.text = datetime.now().strftime("%Y-%m-%d")
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

        # Executive summary slide
        slide_exec = prs.slides.add_slide(prs.slide_layouts[6])
        slide_exec.background.fill.solid()
        slide_exec.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

        tx_title = slide_exec.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.2), Inches(0.8))
        tf_title = tx_title.text_frame
        p = tf_title.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
        p.font.bold = True

        tx_exec = slide_exec.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
        tf_exec = tx_exec.text_frame
        p = tf_exec.paragraphs[0]
        p.text = context["snapshot_text"]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
        p.font.bold = True

        for text, color in [
            (overview["headline"], RGBColor(0xF9, 0xFA, 0xFB)),
            (overview["body"], RGBColor(0xD1, 0xD5, 0xDB)),
            (f"Trend: {history_story['headline']}", RGBColor(0xD1, 0xD5, 0xDB)),
            (f"Run diff: {run_diff['body']}", RGBColor(0xD1, 0xD5, 0xDB)),
            (f"Settings: {' / '.join(settings.get('current_state', [])[:3])}", RGBColor(0xD1, 0xD5, 0xDB)),
            (f"Project flow: {context['flow']}", RGBColor(0x9C, 0xA3, 0xAF)),
            (f"Offset reference: {context['offset_text']}", RGBColor(0x9C, 0xA3, 0xAF)),
            (f"Next checks: {' / '.join(overview.get('next_steps', [])[:3])}", RGBColor(0xD1, 0xD5, 0xDB)),
        ]:
            para = tf_exec.add_paragraph()
            para.text = text
            para.font.size = Pt(14)
            para.font.color.rgb = color

        if issue_spotlight or action_items:
            slide_findings = prs.slides.add_slide(prs.slide_layouts[6])
            slide_findings.background.fill.solid()
            slide_findings.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

            tx_title2 = slide_findings.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf_title2 = tx_title2.text_frame
            p = tf_title2.paragraphs[0]
            p.text = "Critical Findings"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            p.font.bold = True

            tx_findings = slide_findings.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(3.2))
            tf_findings = tx_findings.text_frame
            for idx, finding in enumerate(issue_spotlight[:3]):
                para = tf_findings.paragraphs[0] if idx == 0 else tf_findings.add_paragraph()
                para.text = f"[{finding['status']}] {finding['module']} - {finding['title']}: {finding['evidence']}"
                para.font.size = Pt(15)
                para.font.color.rgb = RGBColor(0xEF, 0x44, 0x44) if finding["status"] == "FAIL" else RGBColor(0xF5, 0x9E, 0x0B)

            tx_actions = slide_findings.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(12), Inches(1.8))
            tf_actions = tx_actions.text_frame
            para = tf_actions.paragraphs[0]
            para.text = "Recommended Actions"
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            for idx, line in enumerate(action_items[:3]):
                para = tf_actions.add_paragraph()
                para.text = f"- {line}"
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        # Summary slide
        if done_results:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            slide2.background.fill.solid()
            slide2.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

            txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf2 = txBox2.text_frame
            p = tf2.paragraphs[0]
            p.text = "QC Run History"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            p.font.bold = True

            # Results as text list
            txBox3 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            tf3 = txBox3.text_frame

            summary = tf3.paragraphs[0]
            summary.text = history_story["body"]
            summary.font.size = Pt(14)
            summary.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

            if run_diff.get("changes"):
                diff_para = tf3.add_paragraph()
                diff_para.text = "Run diff: " + run_diff["body"]
                diff_para.font.size = Pt(14)
                diff_para.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

            for history in history_rows:
                score = float(history["score"]) if history["score"] not in ("", "---") else 0.0
                grade = history["grade"]

                p = tf3.add_paragraph()
                p.text = f"{history['snapshot']}  |  {score:.1f}  |  {grade}  |  {history['completed']}"
                p.font.size = Pt(14)

                if score >= 90:
                    p.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)
                elif score >= 75:
                    p.font.color.rgb = RGBColor(0x3B, 0x82, 0xF6)
                elif score >= 60:
                    p.font.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
                else:
                    p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)

        priority_rows = [row for row in module_rows if row["status"] in ("FAIL", "WARNING")]
        if not priority_rows:
            priority_rows = module_rows[:3]
        if priority_rows:
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            slide3.background.fill.solid()
            slide3.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

            txBox4 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            tf4 = txBox4.text_frame
            p = tf4.paragraphs[0]
            p.text = "Priority Modules"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            p.font.bold = True

            txBox5 = slide3.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.8))
            tf5 = txBox5.text_frame
            for idx, story in enumerate(priority_rows[:3]):
                para = tf5.paragraphs[0] if idx == 0 else tf5.add_paragraph()
                para.text = f"{story['module']} [{story['status']}] - {story['current_reading']}"
                para.font.size = Pt(16 if idx == 0 else 15)
                if story["status"] == "FAIL":
                    para.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
                elif story["status"] == "WARNING":
                    para.font.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
                else:
                    para.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        if settings.get("recommendations"):
            slide_settings = prs.slides.add_slide(prs.slide_layouts[6])
            slide_settings.background.fill.solid()
            slide_settings.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

            tx_title3 = slide_settings.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf_title3 = tx_title3.text_frame
            p = tf_title3.paragraphs[0]
            p.text = "Settings Guidance"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            p.font.bold = True

            tx_settings = slide_settings.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.6))
            tf_settings = tx_settings.text_frame
            p = tf_settings.paragraphs[0]
            p.text = "Current state: " + " / ".join(settings.get("current_state", [])[:4])
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
            for line in settings["recommendations"][:4]:
                para = tf_settings.add_paragraph()
                para.text = f"- {line}"
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        if provenance.get("has_data"):
            slide_prov = prs.slides.add_slide(prs.slide_layouts[6])
            slide_prov.background.fill.solid()
            slide_prov.background.fill.fore_color.rgb = RGBColor(0x0A, 0x0E, 0x17)

            tx_title4 = slide_prov.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf_title4 = tx_title4.text_frame
            p = tf_title4.paragraphs[0]
            p.text = "Provenance"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            p.font.bold = True

            tx_prov = slide_prov.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12), Inches(5.8))
            tf_prov = tx_prov.text_frame
            for idx, line in enumerate(_format_export_provenance_lines(provenance)):
                para = tf_prov.paragraphs[0] if idx == 0 else tf_prov.add_paragraph()
                para.text = line
                para.font.size = Pt(15 if idx == 0 else 14)
                para.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

        prs.save(self._output_path)
